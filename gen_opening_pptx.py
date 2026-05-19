"""
开题报告 PPT 生成脚本
====================

生成《基于对比置信度的领域知识挟持与软引导投机解码》开题报告 PPT。

用法：
    /opt/anaconda3/bin/python3 gen_opening_pptx.py

输出：开题报告_PPT.pptx（共约 23 页）

依赖：python-pptx
图片占位：所有 fig_XX.png 位置先放空白占位框（带 nano banana prompt 文字提示），
用户后续用 nano banana 生成 figures/fig_XX.png 后，重新运行本脚本会自动嵌入。
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# 全局配置
# ============================================================
OUT_PPTX = "开题报告_PPT.pptx"
FIG_DIR = Path("figures")  # nano banana 生成的图放这里：figures/fig_01.png ...

# 16:9 幻灯片尺寸（13.33 × 7.5 inch）
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色（学术风：navy / teal / orange / dark gray）
NAVY   = RGBColor(0x1F, 0x3A, 0x68)
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
GRAY   = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT  = RGBColor(0xF4, 0xF1, 0xDE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)

# 元数据占位（用户后续手动替换）
META = {
    "title": "基于对比置信度的领域知识挟持与软引导投机解码",
    "subtitle": "Contrastive-Confidence Guided Speculative Decoding\n+ Domain Absorption Flywheel",
    "author": "[姓名占位]",
    "advisor": "[导师姓名占位]",
    "school": "[学校 / 学院占位]",
    "major":  "[专业占位]",
    "date":   "2026 年 4 月",
}


# ============================================================
# 工具函数
# ============================================================
def add_blank_slide(prs):
    """添加 16:9 空白幻灯片，背景为白色。"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font_name="Microsoft YaHei"):
    """添加一个文本框，支持多行（用 \n 分隔）。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_color_bar(slide, left, top, width, height, color=NAVY):
    """在幻灯片顶部加一个色条，用于章节标题装饰。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_slide_header(slide, chapter_no, chapter_name, page_no, total_pages=23):
    """统一的页面头部：左上角章节号 + 章节名，右上角页码。"""
    # 顶部色条（厚 8 pt）
    add_color_bar(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.12), NAVY)

    # 左上：章节标题
    add_textbox(slide, Inches(0.4), Inches(0.18), Inches(9), Inches(0.5),
                f"第 {chapter_no} 章  {chapter_name}",
                font_size=14, bold=True, color=NAVY)

    # 右上：页码
    add_textbox(slide, Inches(11.5), Inches(0.18), Inches(1.6), Inches(0.5),
                f"{page_no} / {total_pages}",
                font_size=12, color=GRAY, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, top=Inches(0.75)):
    """页面正文区的大标题（位于头部下方）。"""
    add_textbox(slide, Inches(0.5), top, Inches(12.3), Inches(0.7),
                title, font_size=26, bold=True, color=NAVY)
    # 标题下的细分割线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), top + Inches(0.65),
                                  Inches(12.3), Emu(20000))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_image_or_placeholder(slide, fig_name, left, top, width, height,
                             prompt_text=""):
    """如果 figures/fig_name 存在则插入图片；否则放占位框 + prompt 文字。"""
    fig_path = FIG_DIR / fig_name
    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), left, top, width, height)
        return

    # 占位：带浅灰边框的虚线矩形
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT
    rect.line.color.rgb = GRAY
    rect.line.width = Pt(1)

    add_textbox(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.4), Inches(0.4),
                f"[图片占位：{fig_name}]",
                font_size=12, bold=True, color=NAVY)

    add_textbox(slide, left + Inches(0.2), top + Inches(0.55),
                width - Inches(0.4), height - Inches(0.7),
                f"Nano Banana Prompt：\n{prompt_text}",
                font_size=10, color=GRAY)


def add_table(slide, left, top, width, height, header, rows,
              header_color=NAVY, font_size=12):
    """添加一个简洁的表格。header 是 list[str]，rows 是 list[list[str]]。"""
    n_rows = len(rows) + 1
    n_cols = len(header)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # header
    for j, txt in enumerate(header):
        cell = tbl.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Microsoft YaHei"

    # rows
    for i, row in enumerate(rows, start=1):
        for j, txt in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(txt)
            run.font.size = Pt(font_size - 1)
            run.font.color.rgb = BLACK
            run.font.name = "Microsoft YaHei"
    return tbl


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, line_spacing=1.4, color=BLACK):
    """添加项目符号列表。items 是 list[str]，每条作为一个 bullet。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        # 用三角形符号代替默认 bullet（对中文 PPT 更友好）
        run.text = f"▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Microsoft YaHei"
    return tb


# ============================================================
# 各页生成函数
# ============================================================

def slide_cover(prs):
    """封面页。"""
    s = add_blank_slide(prs)
    # 顶部装饰条
    add_color_bar(s, Emu(0), Emu(0), SLIDE_W, Inches(2.0), NAVY)
    # 中部副条
    add_color_bar(s, Emu(0), Inches(2.0), SLIDE_W, Inches(0.06), TEAL)

    # 主题
    add_textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                "博士学位论文开题报告",
                font_size=20, bold=True, color=WHITE)

    add_textbox(s, Inches(0.6), Inches(0.95), Inches(12), Inches(1.0),
                META["title"],
                font_size=32, bold=True, color=WHITE)

    add_textbox(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
                META["subtitle"].replace("\n", "  "),
                font_size=16, color=LIGHT)

    # 中部留白处放一张总览图
    add_image_or_placeholder(s, "fig_01.png",
                             Inches(2.5), Inches(2.4),
                             Inches(8.3), Inches(3.4),
                             prompt_text="见 nano_banana_prompts.md → fig_01")

    # 底部信息
    add_textbox(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
                f"汇报人：{META['author']}    指导教师：{META['advisor']}",
                font_size=18, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
                f"{META['school']}    {META['major']}",
                font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                META["date"],
                font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


def slide_toc(prs):
    """目录页。"""
    s = add_blank_slide(prs)
    add_color_bar(s, Emu(0), Emu(0), SLIDE_W, Inches(0.12), NAVY)
    add_slide_title(s, "目  录")

    chapters = [
        ("01", "研究背景及现状",   "Background & State of the Art"),
        ("02", "研究目标",         "Research Objectives"),
        ("03", "研究内容",         "Research Content"),
        ("04", "目前进展",         "Current Progress"),
        ("05", "进度安排",         "Schedule & Risk Plan"),
    ]
    top = Inches(2.0)
    for i, (no, name_zh, name_en) in enumerate(chapters):
        # 序号块
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(2.5), top + Inches(i * 0.85),
                                  Inches(0.9), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        add_textbox(s, Inches(2.5), top + Inches(i * 0.85),
                    Inches(0.9), Inches(0.7), no,
                    font_size=24, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 章节名
        add_textbox(s, Inches(3.7), top + Inches(i * 0.85),
                    Inches(7), Inches(0.4), name_zh,
                    font_size=22, bold=True, color=NAVY)
        add_textbox(s, Inches(3.7), top + Inches(i * 0.85 + 0.4),
                    Inches(7), Inches(0.3), name_en,
                    font_size=12, color=GRAY)


# ---------- 第 1 章：研究背景及现状 ----------

def slide_ch1_p1(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 1, "研究背景及现状", page)
    add_slide_title(s, "1.1 为什么这个问题很重要：速度与能力的根本矛盾")

    add_bullet_list(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(4),
                    items=[
                        "32B 通用大模型能力强，但推理慢、部署贵，难以支撑高频实时交互场景",
                        "医疗 / 法律 / 金融等垂直领域真正需要的是“又懂专业、又能快速响应”的系统",
                        "小模型可以被微调成领域专家，但它并不具备大模型那样稳定的通用推理能力",
                        "工业界真正想要的不是“二选一”，而是让大模型保留主体能力、小模型补足领域专知",
                        "投机解码天然提供了这种协作框架，因此它是连接“强能力”和“低时延”的关键技术入口",
                    ],
                    font_size=15)

    add_image_or_placeholder(s, "fig_02.png",
                             Inches(8.3), Inches(1.7),
                             Inches(4.7), Inches(4.5),
                             prompt_text="见 nano_banana_prompts.md → fig_02\n（H200 单卡三模型协同架构图）")

    # 底部一句话总结
    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
                "一句话：垂直领域最需要“大模型 + 小专家”的协同推理，而投机解码本应是这条路的核心基础设施。",
                font_size=15, bold=True, color=ORANGE)


def slide_ch1_p2(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 1, "研究背景及现状", page)
    add_slide_title(s, "1.2 为什么它更重要：垂直领域恰恰是传统 SD 最失效的地方")

    # 左：原理
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(6.2), Inches(0.4),
                "标准 SD 的理想图景：",
                font_size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.5), Inches(2.1), Inches(6.2), Inches(2.5),
                    items=[
                        "小模型先提案：让轻量 Draft 一次生成多个候选 token",
                        "大模型再验收：让 Target 并行验证这些 token 是否可信",
                        "若接受率足够高，就能用少量 Target 调用换来更高吞吐",
                    ],
                    font_size=13)

    # 验收公式
    add_textbox(s, Inches(0.5), Inches(4.0), Inches(6.2), Inches(0.5),
                "P_accept = min(1, P_target(x) / P_draft(x))",
                font_size=16, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)

    # 右：失效数据
    add_textbox(s, Inches(7.0), Inches(1.7), Inches(6), Inches(0.4),
                "现实结果（Surgery, n=200）：",
                font_size=15, bold=True, color=NAVY)
    add_table(s, Inches(7.0), Inches(2.15), Inches(6), Inches(2.0),
              header=["策略", "acc", "tps", "acc_rate"],
              rows=[
                  ["pure_target (32B)", "0.650", "27.3", "1.000"],
                  ["standard_sd",        "0.650", "5.3",  "0.212"],
                  ["hard_override",      "0.650", "17.8", "0.942"],
              ], font_size=12)

    add_textbox(s, Inches(0.5), Inches(5.0), Inches(12.5), Inches(2.0),
                "关键观察：\n"
                "①  standard SD 的 acc 没有提升，但 tps 从 27.3 掉到 5.3\n"
                "②  接受率只有 0.212，说明 Draft 提出的专业 token 大量被拒绝\n"
                "③  hard override 虽然能把速度拉回来，但不会自动带来更高准确率\n"
                "→  这说明问题不是“专家小模型没用”，而是“现有验收机制不会用专家知识”",
                font_size=14, color=GRAY)


def slide_ch1_p3(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 1, "研究背景及现状", page)
    add_slide_title(s, "1.3 为什么这个问题很难：现有三条路线都差最后一环")

    add_table(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(5.4),
              header=["方向", "代表工作", "核心思想", "局限性"],
              rows=[
                  ["投机解码加速",
                   "SpecDec / Medusa / Eagle / SpecInfer",
                   "并行验收 / 多 head 草稿 / 树状解码",
                   "未涉及领域知识；大体量差下接受率仍低"],
                  ["在线分布融合",
                   "Product of Experts / Steering Vec",
                   "logit 域加权乘积 / 激活流注入",
                   "全程注入易过矫正；缺少稀疏门控"],
                 ["领域微调",
                  "LoRA / Adapter / Med-PaLM",
                  "在大模型上微调，让大模型'内化'领域",
                  "训练成本高；通用能力可能塌陷"],
                 ["PEFT 选层 / 自适应 rank",
                  "AdaLoRA / Flexora / GoRA / IGU-LoRA",
                  "梯度 / Hessian / 激活做层重要性排序",
                  "监督锚点为 task loss，与'具体哪一步决策'无关"],
                  ["Token 级加权 SFT",
                   "GIFT / ProFit / ssToken",
                   "按 token entropy / 概率加权 loss",
                   "事件定义粗；与解码事件不耦合"],
                  ["机制可解释性",
                   "Activation Patching / Causal Tracing",
                   "反事实激活替换定位关键层",
                   "工具箱性质，本身非系统方法"],
              ], font_size=12)

    add_textbox(s, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
                "→ 现有工作分别解决了“加速”“注入”或“微调”，但没有把三者真正串成一个闭环故事。",
                font_size=14, bold=True, color=ORANGE)


def slide_ch1_p4(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 1, "研究背景及现状", page)
    add_slide_title(s, "1.4 为什么它很难：我们面对的三个核心技术难点")

    # 三栏：缺陷
    cols = [
        ("难点一：如何识别真正的领域 token",
         "不能只看 Draft 高概率；\n必须区分“通用高频词”\n和“领域增量知识词”。",
         ORANGE),
        ("难点二：如何注入而不破坏通用能力",
         "硬覆盖太粗暴，全程融合太激进；\n必须只在关键位置、\n按需注入专家知识。",
         TEAL),
        ("难点三：如何把在线收益沉淀为离线能力",
         "即使在线引导有效，\n也未回答这些信号能否被\n吸收到 Target 权重中。",
         NAVY),
    ]
    for i, (title, body, color) in enumerate(cols):
        left = Inches(0.4 + i * 4.3)
        # 顶部色条
        add_color_bar(s, left, Inches(1.7), Inches(4.1), Inches(0.5), color)
        add_textbox(s, left, Inches(1.7), Inches(4.1), Inches(0.5),
                    title, font_size=15, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 内容框
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, Inches(2.2),
                                  Inches(4.1), Inches(2.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = LIGHT
        rect.line.color.rgb = color
        add_textbox(s, left + Inches(0.2), Inches(2.3),
                    Inches(3.7), Inches(1.8),
                    body, font_size=13, color=BLACK)

    # 切入点
    add_textbox(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.5),
                "→ 因此，本文的切入点不是继续堆更强 Draft，而是重写“知识如何进入系统”的路径：",
                font_size=18, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.6), Inches(5.1), Inches(12), Inches(2.2),
                    items=[
                        "先解决“怎么借”——在 SD 验收链路中显式引入领域知识探针 ΔP",
                        "再解决“怎么控”——只在 Draft 真有优势且 Target 真不确定时做软引导",
                        "最后解决“怎么吸收”——把 flip 事件转为 PEFT 选层与训练信号",
                        "由此形成'在线引导 → 离线吸收'的完整闭环",
                    ], font_size=14)


# ---------- 第 2 章：研究目标 ----------

def slide_ch2_p1(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 2, "研究目标", page)
    add_slide_title(s, "2.1 我们要回答的三个核心问题")

    add_textbox(s, Inches(0.5), Inches(1.7), Inches(12.5), Inches(0.6),
                "总体目标：",
                font_size=18, bold=True, color=NAVY)
    add_textbox(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.2),
                "围绕“大模型如何高效使用并最终吸收小模型的领域专知”这一主线，\n"
                "构建一个兼顾在线增益、离线沉淀与通用能力保护的双闭环框架。",
                font_size=16, color=BLACK)

    sub_goals = [
        ("问题 A · 怎么让大模型在推理时借到小模型的领域知识",
         "在零训练前提下，重写 SD 验收机制，使 Draft 的领域优势不再被 Target 系统性拒绝。",
         NAVY),
        ("问题 B · 怎么保证这种借用是精准、稀疏、可控的",
         "避免硬覆盖或全程融合对通用流利度的破坏，只在真正需要帮助的位置放大领域信号。",
         TEAL),
        ("问题 C · 怎么把在线增益沉淀为 Target 自身能力",
         "利用 flip 事件指导最小 LoRA 增量，验证领域知识是否能从运行时协同过渡到参数化吸收。",
         ORANGE),
    ]
    for i, (title, body, color) in enumerate(sub_goals):
        top = Inches(3.7 + i * 1.15)
        # 左侧色条
        add_color_bar(s, Inches(0.5), top, Inches(0.18), Inches(1.0), color)
        add_textbox(s, Inches(0.85), top, Inches(12), Inches(0.4),
                    title, font_size=16, bold=True, color=color)
        add_textbox(s, Inches(0.85), top + Inches(0.4), Inches(12), Inches(0.7),
                    body, font_size=13, color=BLACK)


def slide_ch2_p2(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 2, "研究目标", page)
    add_slide_title(s, "2.2 我们做了什么创新：两条创新线，一个统一故事")

    # 左：创新点一
    add_color_bar(s, Inches(0.4), Inches(1.7), Inches(6.1), Inches(0.55), NAVY)
    add_textbox(s, Inches(0.4), Inches(1.7), Inches(6.1), Inches(0.55),
                "创新点 ①  DSSD：先证明“知识可以被借到”",
                font_size=17, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rect_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.4), Inches(2.25),
                                Inches(6.1), Inches(4.7))
    rect_l.fill.solid(); rect_l.fill.fore_color.rgb = LIGHT
    rect_l.line.color.rgb = NAVY

    add_bullet_list(s, Inches(0.5), Inches(2.4), Inches(5.9), Inches(4.4),
                    items=[
                        "提出 ΔP = P_draft − P_base，用 Base 对照显式提取“领域增量知识”",
                        "提出软引导验收机制，把“拒绝专家”改成“按需放大专家 token”",
                        "从 C1 到 C8 建立完整策略谱系，解释什么形式的引导最有效",
                        "在 Surgery 上把 32B Target 准确率从 0.650 提升到 0.700",
                        "说明领域知识在推理期确实可以被大模型借到",
                    ], font_size=13)

    # 右：创新点二
    add_color_bar(s, Inches(6.85), Inches(1.7), Inches(6.1), Inches(0.55), TEAL)
    add_textbox(s, Inches(6.85), Inches(1.7), Inches(6.1), Inches(0.55),
                "创新点 ②  DAF：再证明“知识可以被吸收到”",
                font_size=17, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rect_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(6.85), Inches(2.25),
                                Inches(6.1), Inches(4.7))
    rect_r.fill.solid(); rect_r.fill.fore_color.rgb = LIGHT
    rect_r.line.color.rgb = TEAL

    add_bullet_list(s, Inches(6.95), Inches(2.4), Inches(5.9), Inches(4.4),
                    items=[
                        "把 DSSD 中的 flip 事件定义为“领域知识真实改写决策边界”的证据",
                        "用 flip 事件做 FDLP 选层，让 LoRA 预算投向真正有杠杆的层",
                        "构建 Decode → Place → Train → Re-decode 的 K 轮飞轮",
                        "用 entropy / disagreement / all-token 事件做严格对照，保护创新边界",
                        "目标是把在线收益沉淀为离线能力，而不是永远依赖三模型共跑",
                    ], font_size=13)

    add_textbox(s, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
                "统一故事线：先回答“为什么值得做”，再回答“为什么难”，最后给出“怎么创新地做成它”。",
                font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ---------- 第 3 章：研究内容 ----------

def slide_ch3_p1(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.1 总体技术路线：从“借知识”到“吸知识”")

    add_image_or_placeholder(s, "fig_01.png",
                             Inches(0.7), Inches(1.7),
                             Inches(8.5), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_01\n（DSSD + DAF 协同进化总体技术路线）")

    # 右侧文字说明
    add_textbox(s, Inches(9.5), Inches(1.7), Inches(3.5), Inches(0.4),
                "一句话理解：",
                font_size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(9.5), Inches(2.1), Inches(3.5), Inches(5),
                    items=[
                        "第一步：\n先让 Target 学会在推理时“借”到专家知识",
                        "第二步：\n记录哪些 token 决策真正被领域知识改写",
                        "第三步：\n把这些事件转成 LoRA 该加在哪些层的信号",
                        "第四步：\n让更新后的 Target 再运行，观察是否不再需要这些外援",
                        "最终目标：\n从在线协同走向离线吸收",
                    ], font_size=11, line_spacing=1.2)


def slide_ch3_p2(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.2 创新点一：DSSD 如何解决“知识借不过来”")

    # 左：信号定义与公式
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(6), Inches(0.5),
                "核心回答：不是直接信任 Draft，而是先定位“哪里是真正的领域增量”",
                font_size=15, bold=True, color=NAVY)
    add_textbox(s, Inches(0.5), Inches(2.2), Inches(6), Inches(1.0),
                "ΔP(x) = softmax(logit_draft / T_fixed)[x]\n"
                "          − softmax(logit_base / T_fixed)[x]",
                font_size=14, color=BLACK,
                font_name="Consolas")

    add_textbox(s, Inches(0.5), Inches(3.4), Inches(6), Inches(0.5),
                "只在高价值位置触发：",
                font_size=15, bold=True, color=NAVY)
    add_textbox(s, Inches(0.5), Inches(3.9), Inches(6), Inches(0.7),
                "Condition_Domain ⟺  P_draft(x) > θ_high  ∧  ΔP > τ",
                font_size=14, color=BLACK, font_name="Consolas")

    add_textbox(s, Inches(0.5), Inches(4.8), Inches(6), Inches(0.5),
                "把领域信号写进验收概率：",
                font_size=15, bold=True, color=NAVY)
    add_textbox(s, Inches(0.5), Inches(5.3), Inches(6), Inches(0.7),
                "P'_accept = min(1,  P_t / P_d  +  α · ΔP)",
                font_size=15, bold=True, color=ORANGE,
                font_name="Consolas")

    # 右：直觉示意
    add_image_or_placeholder(s, "fig_02.png",
                             Inches(7.0), Inches(1.7),
                             Inches(6.0), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_02\n（三模型协同 + ΔP 信号流示意）")

    # 关键论断
    add_textbox(s, Inches(0.5), Inches(6.5), Inches(6), Inches(0.7),
                "关键洞察：问题不在于 Draft 不够专业，而在于原始验收公式不会用这份专业性。\n"
                "DSSD 的本质是把“领域增量知识”显式写回验收链路。",
                font_size=12, color=GRAY)


def slide_ch3_p3(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.3 创新点一：我们如何一步步把方法做对")

    add_image_or_placeholder(s, "fig_03.png",
                             Inches(0.5), Inches(1.7),
                             Inches(8.0), Inches(5.3),
                             prompt_text="见 nano_banana_prompts.md → fig_03\n（C1–C9 策略演进树状图）")

    # 右侧关键策略表
    add_textbox(s, Inches(8.8), Inches(1.7), Inches(4.4), Inches(0.4),
                "策略演进的故事：",
                font_size=14, bold=True, color=NAVY)
    add_table(s, Inches(8.8), Inches(2.15), Inches(4.4), Inches(4.5),
              header=["策略", "信号", "best acc"],
              rows=[
                  ["C1", "固定 α，比值域", "0.690"],
                  ["C3", "概率域补贴", "0.660"],
                  ["C4", "Draft 自信门", "0.670"],
                  ["C5", "Target 熵权", "0.685"],
                  ["C6", "C4×C5 步级", "0.690"],
                  ["C8", "token 级门控 ★", "0.700"],
                  ["C9", "二值 + 线性", "0.690"],
              ], font_size=11)

    add_textbox(s, Inches(8.8), Inches(6.7), Inches(4.4), Inches(0.5),
                "★ 结论：门控越精准，领域知识注入越有效",
                font_size=12, bold=True, color=ORANGE)


def slide_ch3_p4(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.4 创新点一：为什么必须是“双信号”而不是“单信号”")

    add_image_or_placeholder(s, "fig_04.png",
                             Inches(0.5), Inches(1.7),
                             Inches(7.5), Inches(4.8),
                             prompt_text="见 nano_banana_prompts.md → fig_04\n（C6 双信号 AND 门控架构图）")

    # 右：公式与机理
    add_textbox(s, Inches(8.3), Inches(1.7), Inches(4.7), Inches(0.5),
                "C6 动态 α 公式：",
                font_size=15, bold=True, color=NAVY)
    add_textbox(s, Inches(8.3), Inches(2.2), Inches(4.7), Inches(1.0),
                "α_t = λ · 𝟙(S_t > τ) · S_t · H_t / H_max",
                font_size=12, bold=True, color=BLACK, font_name="Consolas")

    add_textbox(s, Inches(8.3), Inches(3.3), Inches(4.7), Inches(0.4),
                "为什么要两个信号：",
                font_size=14, bold=True, color=NAVY)
    add_bullet_list(s, Inches(8.3), Inches(3.7), Inches(4.7), Inches(2.5),
                    items=[
                        "只有 Draft 自信：可能会误把通用词也当成领域词",
                        "只有 Target 高熵：可能只是句法犹豫，不一定真缺知识",
                        "两者相乘：只有“专家真懂 + 大模型真不会”时才介入",
                    ], font_size=11, line_spacing=1.3)

    # 底部：实测数据
    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.5), Inches(0.6),
                "实测：C6 说明“精准路由”本身也是一种加速机制；因为注入更准，所以接受率更高、重采样更少。",
                font_size=13, bold=True, color=ORANGE)


def slide_ch3_p5(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.5 创新点二：DAF 回答“能不能把在线增益变成离线能力”")

    add_image_or_placeholder(s, "fig_05.png",
                             Inches(0.5), Inches(1.7),
                             Inches(7.5), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_05\n（DAF 飞轮闭环示意）")

    add_textbox(s, Inches(8.3), Inches(1.7), Inches(4.7), Inches(0.4),
                "飞轮的本质：",
                font_size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(8.3), Inches(2.1), Inches(4.7), Inches(3.0),
                    items=[
                        "先看哪里还要靠外援（flip 事件）",
                        "再看这些外援应该被写进哪几层参数",
                        "然后做最小 LoRA 更新",
                        "最后检查更新后的 Target 是否还需要同类外援",
                    ], font_size=12)

    add_textbox(s, Inches(8.3), Inches(5.3), Inches(4.7), Inches(0.4),
                "为什么是“飞轮”而不是“一次微调”：",
                font_size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(8.3), Inches(5.7), Inches(4.7), Inches(2.0),
                    items=[
                        "因为每一轮都会改变下一轮的 flip 分布",
                        "如果 flip 持续下降，说明知识正在被吸收",
                        "如果通用域退化，就立刻停下并回退",
                    ], font_size=11, line_spacing=1.2)


def slide_ch3_p6(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.6 创新点二：真正的创新锚点是什么")

    # 左：FDLP 流程
    add_image_or_placeholder(s, "fig_06.png",
                             Inches(0.5), Inches(1.7),
                             Inches(6.0), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_06\n（FDLP 选层算法流程）")

    # 右：双重角色
    add_image_or_placeholder(s, "fig_07.png",
                             Inches(6.8), Inches(1.7),
                             Inches(6.2), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_07\n（flip 事件双重角色示意）")

    add_textbox(s, Inches(0.5), Inches(6.85), Inches(12.5), Inches(0.5),
                "核心创新：不是“又做了一个 LoRA”，而是第一次把 speculative decoding 内生事件同时用作训练监督与停止信号。",
                font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_ch3_p7(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 3, "研究内容", page)
    add_slide_title(s, "3.7 如何证明这些创新不是“讲故事”")

    add_table(s, Inches(0.4), Inches(1.7), Inches(12.5), Inches(5.3),
              header=["阶段", "目标", "代表实验", "通过标准"],
              rows=[
                  ["Phase 0",
                   "现象统计 + 事件必要性",
                   "FDLP vs all-token / entropy / disagreement 选层",
                   "Flip 选层显著优于其他事件"],
                  ["Phase 1",
                   "DAF 飞轮主实验",
                   "F0=Target / F1=DSSD / F2-F4=飞轮 K=1,2,3",
                   "F4 acc ≈ F1，tps ≈ F0"],
                  ["Phase 2",
                   "事件对照飞轮",
                   "G1=flip / G2=entropy / G3=disagreement / G4=all-token",
                   "仅 G1 收敛、G2-G4 无效"],
                  ["Phase 3",
                   "强基线对照（固定预算）",
                   "L1=全层 / L2=DAF / L3=随机 / L8=AdaLoRA",
                   "L2 与 L8 持平，显著优于 L1/L3"],
                  ["Phase 4 (可选)",
                   "飞轮收敛律拟合",
                   "ρ_k = 1 − F̄^(k)/F̄^(0) 拟合曲线",
                   "给出 wall-clock 经验律"],
              ], font_size=11)

    add_textbox(s, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
                "→ 关键不是只做出一个有效结果，而是证明：只有我们的事件定义与闭环设计，才能稳定得到这个结果。",
                font_size=13, color=GRAY)


# ---------- 第 4 章：目前进展 ----------

def slide_ch4_p1(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 4, "目前进展", page)
    add_slide_title(s, "4.1 已搭建：H200 单卡三模型内聚架构")

    add_image_or_placeholder(s, "fig_02.png",
                             Inches(0.5), Inches(1.7),
                             Inches(7.0), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_02\n（H200 单卡三模型架构）")

    add_textbox(s, Inches(7.8), Inches(1.7), Inches(5.2), Inches(0.4),
                "工程要点：",
                font_size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(7.8), Inches(2.1), Inches(5.2), Inches(5),
                    items=[
                        "H200 141 GB：32B + 3B + 3B 同卡共驻",
                        "全局 bfloat16，device_map=cuda:0 严格无跨卡",
                        "禁用 generate()，手写 model.forward AR loop",
                        "StaticCache 预分配 KV Buffer",
                        "Shadow Sync：Base 跟随 Draft 推进，Lazy LM Head",
                        "TelemetryLogger 记录全策略遥测",
                    ], font_size=12, line_spacing=1.3)


def slide_ch4_p2(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 4, "目前进展", page)
    add_slide_title(s, "4.2 已完成：Draft 模型微调 + Pareto 前沿")

    # 左：微调结果
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
                "Draft-Surgery 微调结果（n=249）：",
                font_size=14, bold=True, color=NAVY)
    add_table(s, Inches(0.5), Inches(2.15), Inches(6), Inches(2.4),
              header=["模型", "Surgery acc", "Δ vs Base"],
              rows=[
                  ["Pure Target 32B",       "0.650", "—"],
                  ["Base 3B（未微调）",     "0.494", "—"],
                  ["Draft-Surgery 3B (ckpt-1676)", "0.574", "+0.080"],
              ], font_size=11)

    add_textbox(s, Inches(0.5), Inches(4.7), Inches(6), Inches(0.5),
                "数据：MedMCQA Surgery 22K + 25% Alpaca 锚点",
                font_size=12, color=GRAY)
    add_textbox(s, Inches(0.5), Inches(5.1), Inches(6), Inches(0.5),
                "训练：LLaMA-Factory FFT，5 epochs，lr=3e-6",
                font_size=12, color=GRAY)
    add_textbox(s, Inches(0.5), Inches(5.5), Inches(6), Inches(0.5),
                "格式锚点：output 固定为 \"Final answer: X\"",
                font_size=12, color=GRAY)

    # 右：Pareto 热图
    add_image_or_placeholder(s, "fig_08.png",
                             Inches(7.0), Inches(1.7),
                             Inches(6.0), Inches(5.0),
                             prompt_text="见 nano_banana_prompts.md → fig_08\n（acc vs tps Pareto 前沿散点图）")


def slide_ch4_p3(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 4, "目前进展", page)
    add_slide_title(s, "4.3 已完成：DSSD C1–C12 完整结果（Surgery, n=200）")

    add_table(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(5.0),
              header=["策略", "α/λ", "acc", "tps", "acc_rate", "speedup", "亮点"],
              rows=[
                  ["pure_target", "—",    "0.650", "27.3", "1.000", "1.00×", "基线"],
                  ["standard_sd", "—",    "0.650", "5.3",  "0.212", "0.19×", "无加速"],
                  ["C1",          "0.10", "0.690", "6.9",  "0.334", "0.25×", "+4pt"],
                  ["C1",          "1.50", "0.660", "17.7", "0.935", "0.65×", "高速点"],
                  ["C5",          "5.00", "0.685", "8.2",  "0.367", "0.30×", "+3.5pt"],
                  ["C6",          "50",   "0.690", "11.5", "0.600", "0.42×", "Pareto 突破"],
                  ["C8",          "20",   "0.700", "12.2", "0.591", "0.44×", "★ 全局最优"],
                  ["C9",          "100",  "0.690", "15.7", "0.768", "0.57×", "tps 最高"],
                  ["C12",         "20",   "0.665", "16.6", "0.867", "0.61×", "logit 域"],
              ], font_size=11)

    add_textbox(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7),
                "→ 核心结论：在不修改任何模型权重前提下，Surgery acc 0.650 → 0.700（+5%），\n"
                "    且 C9 在 tps=15.7 时仍维持 +4pt acc，证明软引导能'既准且快'。",
                font_size=14, bold=True, color=ORANGE)


def slide_ch4_p4(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 4, "目前进展", page)
    add_slide_title(s, "4.4 已完成：Target 熵分布三科对比探针")

    add_image_or_placeholder(s, "fig_09.png",
                             Inches(0.5), Inches(1.7),
                             Inches(7.5), Inches(4.8),
                             prompt_text="见 nano_banana_prompts.md → fig_09\n（三科目 Target 熵直方图）")

    # 右：核心数据
    add_textbox(s, Inches(8.3), Inches(1.7), Inches(4.7), Inches(0.4),
                "三科核心数字：",
                font_size=14, bold=True, color=NAVY)
    add_table(s, Inches(8.3), Inches(2.1), Inches(4.7), Inches(2.5),
              header=["科目", "acc", "p 值"],
              rows=[
                  ["Surgery",       "59%", "0.0008 ✓"],
                  ["Pharmacology",  "76%", "0.1218 ✗"],
                  ["Anatomy",       "74%", "0.0004 ✓"],
              ], font_size=11)

    add_textbox(s, Inches(8.3), Inches(4.7), Inches(4.7), Inches(2.5),
                "结论：\n"
                "•  错题熵显著高于对题（Surgery / Anatomy）\n"
                "•  Pharmacology 不显著 → C5 自动退化\n"
                "•  领域词位置熵 +57%（H=0.611 vs 0.390）\n"
                "→ 直接支撑 C5/C6 设计合理性",
                font_size=11, color=BLACK)

    add_textbox(s, Inches(0.5), Inches(6.7), Inches(7.5), Inches(0.5),
                "（即'领域熵探测器' — Target 在领域盲区处熵持续偏高）",
                font_size=12, color=GRAY)


def slide_ch4_p5(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 4, "目前进展", page)
    add_slide_title(s, "4.5 已完成：Draft 延续词探针（实验 D）")

    # 左：实验设计
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.4),
                "实验设计：",
                font_size=14, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.5), Inches(2.1), Inches(5.5), Inches(2.5),
                    items=[
                        "在 Target 高熵步（H > μ+σ）",
                        "用 Draft 从同上下文贪婪续写 3 token",
                        "统计 Draft 想输出的内容",
                        "100 题 Surgery，41 错题，355 探针事件",
                    ], font_size=12)

    # 右：Top-10 表
    add_textbox(s, Inches(6.3), Inches(1.7), Inches(6.7), Inches(0.4),
                "Draft 延续词 Top-10（Target 困惑时 Draft 想说什么）：",
                font_size=12, bold=True, color=NAVY)
    add_table(s, Inches(6.3), Inches(2.1), Inches(6.7), Inches(4.5),
              header=["排名", "Draft 续词", "频次", "类别"],
              rows=[
                  ["1",  "Final answer:",        "17", "推理噪音"],
                  ["4",  "condylar",             "4",  "★ 颌面术语"],
                  ["7",  "displacement of the",  "2",  "★ 力学"],
                  ["8",  "avoids the risks",     "2",  "★ 情境"],
                  ["9",  "maxilla.",             "2",  "★ 解剖"],
                  ["10", "ductal stricture",     "2",  "★ 病理"],
              ], font_size=10)

    add_textbox(s, Inches(0.5), Inches(4.7), Inches(5.5), Inches(2.5),
                "因果验证（实验 D-2）：\n"
                "•  阈值 τ=2.0 时硬替换 Target token\n"
                "•  仅 1.6% 替换率 →  acc 67.0%（+2.5pt）\n"
                "•  低阈值（τ=0.5）硬替换 → acc 60.5%\n"
                "→ 论证'软引导优于硬替换'",
                font_size=11, color=BLACK)

    add_textbox(s, Inches(0.5), Inches(6.85), Inches(12.5), Inches(0.5),
                "结论：高熵 + Draft 续词中富含外科 / 颌面专业术语，flip 事件具有真实领域语义。",
                font_size=12, bold=True, color=ORANGE)


# ---------- 第 5 章：进度安排 ----------

def slide_ch5_p1(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 5, "进度安排", page)
    add_slide_title(s, "5.1 总体进度甘特图（2026.05 – 2027.06）")

    add_image_or_placeholder(s, "fig_10.png",
                             Inches(0.5), Inches(1.7),
                             Inches(8.3), Inches(5.3),
                             prompt_text="见 nano_banana_prompts.md → fig_10\n（剩余 14 个月研究进度甘特图）")

    add_textbox(s, Inches(9.0), Inches(1.7), Inches(4.0), Inches(0.4),
                "三大里程碑：",
                font_size=14, bold=True, color=NAVY)
    add_bullet_list(s, Inches(9.0), Inches(2.1), Inches(4.0), Inches(2.5),
                    items=[
                        "M1（2026.08）\nDSSD 论文投稿",
                        "M2（2026.12）\n飞轮机制验证",
                        "M3（2027.05）\n论文初稿完成",
                    ], font_size=12, line_spacing=1.3)

    add_textbox(s, Inches(9.0), Inches(5.0), Inches(4.0), Inches(0.4),
                "每月评审节点：",
                font_size=14, bold=True, color=NAVY)
    add_bullet_list(s, Inches(9.0), Inches(5.4), Inches(4.0), Inches(2.0),
                    items=[
                        "月度组会汇报",
                        "季度博士进展评估",
                        "年度博士资格确认",
                    ], font_size=11, line_spacing=1.2)


def slide_ch5_p2(prs, page):
    s = add_blank_slide(prs)
    add_slide_header(s, 5, "进度安排", page)
    add_slide_title(s, "5.2 风险识别与备选方案")

    risks = [
        ("风险一：飞轮 1 轮即收敛（K* = 1）",
         "→  方案降级为'FDLP 单轮 + 收敛判据'，仍是一篇有效成果",
         ORANGE),
        ("风险二：飞轮收敛到 Target = Draft 复制品",
         "→  通用域 MMLU 退化 > 1% 时强制回退至 v_{k-1}，并停止飞轮",
         TEAL),
        ("风险三：flip 热点跨轮漂移（Jaccard < 0.3）",
         "→  改报告'飞轮沿轨迹拓宽吸收带'，作为新发现写入论文",
         NAVY),
        ("风险四：事件对照飞轮（G2/G3）也能收敛",
         "→  flip 不可替代性主张收缩为'工程消融与经验总结'",
         ORANGE),
        ("风险五：32B LoRA 训练显存爆炸",
         "→  用 8-bit AdamW + gradient checkpointing；必要时降至 7B Target",
         GRAY),
    ]
    top = Inches(1.7)
    for i, (title, body, color) in enumerate(risks):
        # 左侧色条
        add_color_bar(s, Inches(0.4), top + Inches(i * 1.0),
                      Inches(0.18), Inches(0.85), color)
        add_textbox(s, Inches(0.75), top + Inches(i * 1.0),
                    Inches(12.2), Inches(0.4),
                    title, font_size=15, bold=True, color=color)
        add_textbox(s, Inches(0.75), top + Inches(i * 1.0 + 0.42),
                    Inches(12.2), Inches(0.5),
                    body, font_size=13, color=BLACK)


# ---------- 致谢页 ----------

def slide_thanks(prs):
    s = add_blank_slide(prs)
    add_color_bar(s, Emu(0), Emu(0), SLIDE_W, Inches(2.5), NAVY)
    add_color_bar(s, Emu(0), Inches(2.5), SLIDE_W, Inches(0.06), TEAL)

    add_textbox(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.2),
                "Thanks for Your Attention",
                font_size=44, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
                "感谢各位老师的聆听，恳请批评指正",
                font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6),
                META["title"],
                font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                META["subtitle"].replace("\n", "  "),
                font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
                f"汇报人：{META['author']}    指导教师：{META['advisor']}",
                font_size=16, color=BLACK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
                f"{META['school']}    {META['major']}    {META['date']}",
                font_size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 主流程
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 封面 + 目录
    slide_cover(prs)            # p1
    slide_toc(prs)              # p2

    # 第 1 章
    slide_ch1_p1(prs, page=3)
    slide_ch1_p2(prs, page=4)
    slide_ch1_p3(prs, page=5)
    slide_ch1_p4(prs, page=6)

    # 第 2 章
    slide_ch2_p1(prs, page=7)
    slide_ch2_p2(prs, page=8)

    # 第 3 章
    slide_ch3_p1(prs, page=9)
    slide_ch3_p2(prs, page=10)
    slide_ch3_p3(prs, page=11)
    slide_ch3_p4(prs, page=12)
    slide_ch3_p5(prs, page=13)
    slide_ch3_p6(prs, page=14)
    slide_ch3_p7(prs, page=15)

    # 第 4 章
    slide_ch4_p1(prs, page=16)
    slide_ch4_p2(prs, page=17)
    slide_ch4_p3(prs, page=18)
    slide_ch4_p4(prs, page=19)
    slide_ch4_p5(prs, page=20)

    # 第 5 章
    slide_ch5_p1(prs, page=21)
    slide_ch5_p2(prs, page=22)

    # 致谢
    slide_thanks(prs)           # p23

    prs.save(OUT_PPTX)
    print(f"[OK] 已生成 {OUT_PPTX}（共 {len(prs.slides)} 页）")


if __name__ == "__main__":
    main()
